# python example that makes 2 png time series plot
######################################################
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd

here = "office/ppt"

# seed
np.random.seed(123)

# ts example
data = pd.DataFrame(
    np.random.randn(10, 4),
    index=pd.date_range("1/1/2000", periods=10),
    columns=list("ABCD"),
)

# plot and save as png
data.plot()
plt.savefig(f"{here}/ts.png")

# another ts example
data = data * 20
data.plot()
plt.savefig(f"{here}/ts2.png")


# python example that puts 2 png on dir to a ppt
######################################################
from pathlib import Path

import pptx.util
from pptx import Presentation


def add_a_png_slide(prs, png: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout 6 is blank

    # add pic in the middle of the slide
    # then resize
    # total height is 7.5 inches
    # title height is 0.5 inches
    # total width is 10 inches
    left = pptx.util.Inches(2)
    top = pptx.util.Inches(1)
    pic_size = 6
    pic = slide.shapes.add_picture(png, left, top)
    pic.width = pptx.util.Inches(pic_size)
    pic.height = pptx.util.Inches(pic_size)


def add_a_2pngs_slide(prs, png1, png2):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # add pic left and right then resize
    top = pptx.util.Inches(1)
    pic_size = 4
    left_margin = 0.5

    left = pptx.util.Inches(left_margin)
    pic = slide.shapes.add_picture(png1, left, top)
    pic.width = pic.height = pptx.util.Inches(pic_size)

    top = pptx.util.Inches(left_margin + pic_size + 1)
    pic.width = pic.height = pptx.util.Inches(pic_size)


# get pngs
pngs = [i.as_posix() for i in Path(here).glob("*.png")]

# delete ppt if exists
ppt = Path(f"{here}/ts.pptx")
if ppt.exists():
    ppt.unlink()

# make ppt
prs = Presentation()

for png in pngs:
    add_a_png_slide(prs, png)

prs.save(f"{here}/ts.pptx")


# a title page
######################################################
# make ppt
prs = Presentation()

# title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Office"
subtitle.text = "subtitle"

prs.save(f"{here}/title.pptx")


# functionalise
def add_title(prs, title_text, subtitle_text):
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text
    subtitle.text = subtitle_text


# orchestrate
###########################

# delete ppt if exists
ppt = Path(f"{here}/orchestrate.pptx")
if ppt.exists():
    ppt.unlink()

# make ppt
prs = Presentation()

# a title slide and 2 pngs
add_title(prs, "Office", "subtitle")
add_pngs(prs, pngs)

# another title slide and 2 pngs
add_title(prs, "Office2", "subtitle2")
add_pngs(prs, pngs)

prs.save(f"{here}/orchestrate.pptx")
